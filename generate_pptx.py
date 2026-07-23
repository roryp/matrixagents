"""Generate a widescreen PPTX presentation for MatrixAgents showcase."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import os

DOCS = os.path.join(os.path.dirname(__file__), "docs")

# Widescreen 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# Colors
BG_COLOR = RGBColor(0x0D, 0x11, 0x17)       # dark background
TITLE_COLOR = RGBColor(0x58, 0xA6, 0xFF)     # blue accent
TEXT_COLOR = RGBColor(0xE6, 0xED, 0xF3)       # light text
SUBTITLE_COLOR = RGBColor(0x8B, 0x94, 0x9E)   # muted grey
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)     # green accent
TABLE_HEADER_BG = RGBColor(0x16, 0x1B, 0x22)
TABLE_ROW_BG = RGBColor(0x0D, 0x11, 0x17)


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_image_centered(slide, img_path, top, max_width=None, max_height=None, crop_top=0):
    """Add image centered horizontally, fitting within max dimensions."""
    from PIL import Image
    if not os.path.exists(img_path):
        return
    img = Image.open(img_path)
    source = img_path
    if crop_top:
        img = img.crop((0, crop_top, img.width, img.height))
        buffer = BytesIO()
        img.save(buffer, format="PNG")
        buffer.seek(0)
        source = buffer
    img_w, img_h = img.size
    dpi = 96
    emu_w = int(img_w / dpi * 914400)
    emu_h = int(img_h / dpi * 914400)

    if max_width is None:
        max_width = SLIDE_W - Inches(1)
    if max_height is None:
        max_height = SLIDE_H - top - Inches(0.5)

    scale = min(max_width / emu_w, max_height / emu_h, 1.0)
    final_w = int(emu_w * scale)
    final_h = int(emu_h * scale)
    left = int((SLIDE_W - final_w) / 2)
    slide.shapes.add_picture(source, left, int(top), final_w, final_h)


def add_panel(slide, left, top, width, height, border_color=TITLE_COLOR):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = TABLE_HEADER_BG
    panel.line.color.rgb = border_color
    panel.line.width = Pt(1.5)
    return panel


def add_dashboard_focus(slide, img_path, top):
    """Show the live dashboard overview without truncating a pattern section."""
    from PIL import Image

    image = Image.open(img_path)
    dashboard_focus = image.crop((0, 0, image.width, min(440, image.height)))
    buffer = BytesIO()
    dashboard_focus.save(buffer, format="PNG")
    buffer.seek(0)

    panel_width = Inches(11.8)
    panel_height = Inches(4.15)
    image_width, image_height = dashboard_focus.size
    scale = min(panel_width / image_width, panel_height / image_height)
    final_width = int(image_width * scale)
    final_height = int(image_height * scale)
    left = int((SLIDE_W - final_width) / 2)
    image_top = int(top + (panel_height - final_height) / 2)
    slide.shapes.add_picture(buffer, left, image_top, final_width, final_height)

    category_panels = [
        (Inches(0.65), "WORKFLOW · 5", "Sequential · Parallel · Mapper · Loop · Conditional", TITLE_COLOR),
        (Inches(4.72), "AGENTIC · 2", "Supervisor · Human-in-the-Loop", RGBColor(0xC0, 0x84, 0xFC)),
        (Inches(8.79), "PLANNING · 4", "GOAP · P2P · Debate · Voting", RGBColor(0xFB, 0x92, 0x3C)),
    ]
    for panel_left, heading, details, color in category_panels:
        add_panel(slide, panel_left, Inches(5.55), Inches(3.88), Inches(1.05), color)
        add_textbox(slide, panel_left + Inches(0.18), Inches(5.72), Inches(3.52), Inches(0.28),
                    heading, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, panel_left + Inches(0.18), Inches(6.12), Inches(3.52), Inches(0.3),
                    details, font_size=10, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)


# ── Slide 1: Title ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
            "AI Agents", font_size=54, color=TITLE_COLOR, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
            "LangChain4j Agentic Patterns Showcase", font_size=28,
            color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
            "11 agentic patterns with real-time visualization using D3.js and WebSocket streaming",
            font_size=17, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
            "Quarkus 3.30  •  LangChain4j 1.18  •  React 18  •  D3.js  •  Azure OpenAI",
            font_size=16, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

# ── Slide 2: Screenshot ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "Live Dashboard", font_size=36, color=TITLE_COLOR, bold=True,
            alignment=PP_ALIGN.CENTER)
add_dashboard_focus(slide, os.path.join(DOCS, "screenshot.png"), Inches(1.2))

# ── Slide 3: Architecture ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "Azure Deployment Architecture", font_size=36, color=TITLE_COLOR, bold=True,
            alignment=PP_ALIGN.CENTER)
add_image_centered(slide, os.path.join(DOCS, "architecture.png"), Inches(1.2),
                   max_width=Inches(12.2), max_height=Inches(5.9), crop_top=125)

# ── Slide 4: Patterns Overview ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "11 Agentic Patterns Overview", font_size=36, color=TITLE_COLOR, bold=True,
            alignment=PP_ALIGN.CENTER)
add_image_centered(slide, os.path.join(DOCS, "patterns-overview.png"), Inches(1.3))

# ── Pattern slides ──────────────────────────────────────────────
patterns = [
    {
        "title": "1. Sequential Workflow",
        "subtitle": "Chain Orchestration",
        "image": "pattern-sequential.png",
        "bullets": [
            "Agents run one after another, like an assembly line",
            "Each step depends on the previous step's output",
            "Example: CreativeWriter → AudienceEditor → StyleEditor",
            "API: AgenticServices.sequenceBuilder()",
        ],
    },
    {
        "title": "2. Parallel Workflow",
        "subtitle": "Fan-out Orchestration",
        "image": "pattern-parallel.png",
        "bullets": [
            "Multiple agents run simultaneously, results combined",
            "Get diverse perspectives quickly",
            "Example: FoodExpert + MovieExpert → Combiner",
            "API: @Parallel annotation with createAgenticSystem()",
        ],
    },
    {
        "title": "3. Parallel Mapper",
        "subtitle": "Concurrent Map-Reduce",
        "image": "pattern-parallel-mapper.png",
        "layout": "portrait",
        "bullets": [
            "Creates one worker instance per collection item",
            "Processes the full batch concurrently while preserving order",
            "Example: Review batch → ReviewAnalyzer workers → Report",
            "API: AgenticServices.parallelMapperBuilder()",
        ],
    },
    {
        "title": "4. Loop Workflow",
        "subtitle": "Cycle Orchestration",
        "image": "pattern-loop.png",
        "bullets": [
            "Iterative refinement until quality threshold is met",
            "Quality matters more than speed",
            "Example: Writer → Scorer → Editor (repeat until score ≥ 0.8)",
            "API: AgenticServices.loopBuilder()",
        ],
    },
    {
        "title": "5. Conditional Routing",
        "subtitle": "Branch Orchestration",
        "image": "pattern-conditional.png",
        "bullets": [
            "Routes to different specialist agents based on input",
            "Different inputs need different expertise",
            "Example: CategoryRouter → Medical / Legal / Technical Expert",
            "API: @Conditional annotation with createAgenticSystem()",
        ],
    },
    {
        "title": "6. Supervisor Agent",
        "subtitle": "Star Orchestration",
        "image": "pattern-supervisor.png",
        "bullets": [
            "A 'boss' agent plans and delegates to worker agents",
            "LLM intelligence decides task decomposition",
            "Example: BankSupervisor → Withdraw / Credit / Exchange Agent",
            "API: AgenticServices.supervisorBuilder()",
        ],
    },
    {
        "title": "7. Human-in-the-Loop",
        "subtitle": "Gated Orchestration",
        "image": "pattern-humaninloop.png",
        "bullets": [
            "Pauses execution for human input or approval",
            "High-stakes decisions, compliance requirements",
            "Example: ZodiacExtractor → Human Input → HoroscopeAgent",
            "API: AgenticServices.agentBuilder() with WebSocket events",
        ],
    },
    {
        "title": "8. GOAP — Goal-Oriented Action Planning",
        "subtitle": "DAG — Custom Planner",
        "image": "pattern-goap.png",
        "bullets": [
            "Finds optimal sequence of agents to reach a goal",
            "Like GPS finding the shortest route",
            "Example: CityParser → (Distances + Attractions) → Itinerary",
            "API: plannerBuilder() + GoalOrientedPlanner",
        ],
    },
    {
        "title": "9. P2P — Peer-to-Peer",
        "subtitle": "Mesh — Custom Planner",
        "image": "pattern-p2p.png",
        "bullets": [
            "Agents collaborate as equals without a central controller",
            "Emergent collaboration for research & brainstorming",
            "Example: Literature → Hypothesis → Critic → Validation → Scorer",
            "API: plannerBuilder() + P2PPlanner (threshold ≥ 0.75)",
        ],
    },
    {
        "title": "10. Debate",
        "subtitle": "Multi-Round Panel — Custom Planner",
        "image": "pattern-debate.png",
        "layout": "portrait",
        "bullets": [
            "Opposing perspectives refine arguments in parallel rounds",
            "Shared debate context lets agents respond to prior positions",
            "Example: Proponent + Skeptic + Pragmatist → Judge",
            "API: plannerBuilder() + DebatePlanner",
        ],
    },
    {
        "title": "11. Voting",
        "subtitle": "Ensemble Decision — Custom Planner",
        "image": "pattern-voting.png",
        "layout": "portrait",
        "bullets": [
            "Independent specialists cast comparable categorical ballots",
            "A voting strategy aggregates one collective decision",
            "Example: Growth + Value + Risk analysts → Majority",
            "API: plannerBuilder() + VotingPlanner",
        ],
    },
]

for pat in patterns:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                pat["title"], font_size=32, color=TITLE_COLOR, bold=True)
    # Subtitle
    add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.5),
                pat["subtitle"], font_size=16, color=SUBTITLE_COLOR)

    # Use the full slide width for topology diagrams so edge labels remain legible.
    img_path = os.path.join(DOCS, pat["image"])
    if os.path.exists(img_path):
        from PIL import Image
        img = Image.open(img_path)
        crop_top = pat.get("crop_top", 115)
        img = img.crop((0, crop_top, img.width, img.height))
        image_buffer = BytesIO()
        img.save(image_buffer, format="PNG")
        image_buffer.seek(0)
        img_w, img_h = img.size
        dpi = 96
        emu_w = int(img_w / dpi * 914400)
        emu_h = int(img_h / dpi * 914400)
        portrait_layout = pat.get("layout") == "portrait"
        max_w = Inches(6.7 if portrait_layout else 12.1)
        max_h = Inches(5.85 if portrait_layout else 3.85)
        scale = min(max_w / emu_w, max_h / emu_h)
        final_w = int(emu_w * scale)
        final_h = int(emu_h * scale)
        if portrait_layout:
            left = Inches(0.6) + int((max_w - final_w) / 2)
        else:
            left = int((SLIDE_W - final_w) / 2)
        top = Inches(1.35) + int((max_h - final_h) / 2)
        slide.shapes.add_picture(image_buffer, left, top, final_w, final_h)

        if portrait_layout:
            bullets = "\n\n".join(f"•  {bullet}" for bullet in pat["bullets"])
            add_textbox(slide, Inches(7.45), Inches(1.7), Inches(5.2), Inches(5.1),
                bullets, font_size=18, color=TEXT_COLOR)
        else:
            # Keep supporting points in two compact columns below wide diagrams.
            left_bullets = "\n".join(f"•  {b}" for b in pat["bullets"][:2])
            right_bullets = "\n".join(f"•  {b}" for b in pat["bullets"][2:])
            add_textbox(slide, Inches(0.7), Inches(5.35), Inches(5.9), Inches(1.5),
                left_bullets, font_size=13, color=TEXT_COLOR)
            add_textbox(slide, Inches(6.75), Inches(5.35), Inches(5.9), Inches(1.5),
                right_bullets, font_size=13, color=TEXT_COLOR)

# ── Slide: Choosing the Right Pattern ───────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "Choosing the Right Pattern", font_size=36, color=TITLE_COLOR, bold=True,
            alignment=PP_ALIGN.CENTER)

table_data = [
    ("Situation", "Recommended Pattern"),
    ("Simple pipeline with clear steps", "Sequential"),
    ("Need multiple perspectives fast", "Parallel"),
    ("Apply one operation to many items", "Parallel Mapper"),
    ("Quality is critical, time isn't", "Loop"),
    ("Different inputs need different handling", "Conditional"),
    ("Complex task, unclear how to break down", "Supervisor"),
    ("Need human approval or input", "Human-in-the-Loop"),
    ("Many dependencies, need optimal path", "GOAP"),
    ("Creative / brainstorming collaboration", "P2P"),
    ("Need competing arguments examined", "Debate"),
    ("Need a robust categorical decision", "Voting"),
]

rows = len(table_data)
cols = 2
tbl_w = Inches(10)
tbl_h = Inches(5.4)
left = int((SLIDE_W - tbl_w) / 2)
top = Inches(1.5)
table_shape = slide.shapes.add_table(rows, cols, left, top, tbl_w, tbl_h)
table = table_shape.table
table.columns[0].width = Inches(6.5)
table.columns[1].width = Inches(3.5)

for r, (col0, col1) in enumerate(table_data):
    for c, val in enumerate((col0, col1)):
        cell = table.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = "Segoe UI"
        if r == 0:
            p.font.bold = True
            p.font.color.rgb = TITLE_COLOR
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        else:
            p.font.color.rgb = TEXT_COLOR
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_BG

# ── Slide: Tech Stack ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "Tech Stack", font_size=36, color=TITLE_COLOR, bold=True,
            alignment=PP_ALIGN.CENTER)

backend_text = (
    "•  Java 21 with Virtual Threads\n"
    "•  Quarkus 3.30\n"
    "•  LangChain4j 1.18.0 (Core)\n"
    "•  LangChain4j Agentic 1.18.0-beta28\n"
    "•  Agentic Patterns 1.18.0-beta28\n"
    "•  Azure OpenAI integration\n"
    "•  Native Quarkus WebSockets"
)
add_panel(slide, Inches(0.7), Inches(1.45), Inches(5.8), Inches(3.85))
add_textbox(slide, Inches(1.05), Inches(1.8), Inches(5.1), Inches(0.5),
            "Backend + AI", font_size=23, color=TITLE_COLOR, bold=True)
add_textbox(slide, Inches(1.05), Inches(2.45), Inches(5.05), Inches(3.45),
            backend_text, font_size=17, color=TEXT_COLOR)

frontend_text = (
    "•  React 18 with TypeScript\n"
    "•  Vite 5 build tool\n"
    "•  D3.js for visualizations\n"
    "•  Tailwind CSS for styling\n"
    "•  React Router for navigation"
)
add_panel(slide, Inches(6.83), Inches(1.45), Inches(5.8), Inches(3.85), ACCENT_GREEN)
add_textbox(slide, Inches(7.18), Inches(1.8), Inches(5.1), Inches(0.5),
            "Frontend + Visualization", font_size=23, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.18), Inches(2.45), Inches(5.05), Inches(3.45),
            frontend_text, font_size=17, color=TEXT_COLOR)

stack_flow = [
    (Inches(0.95), "Azure OpenAI\nreasoning", TITLE_COLOR),
    (Inches(5.12), "Quarkus + LangChain4j\norchestration", ACCENT_GREEN),
    (Inches(9.28), "React + D3\nlive visualization", RGBColor(0xC0, 0x84, 0xFC)),
]
for left, label, color in stack_flow:
    flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, Inches(5.75), Inches(3.1), Inches(0.95))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = TABLE_HEADER_BG
    flow_box.line.color.rgb = color
    add_textbox(slide, left, Inches(5.93), Inches(3.1), Inches(0.58), label,
                font_size=14, color=TEXT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

for left in (Inches(4.32), Inches(8.48)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                   left, Inches(5.98), Inches(0.55), Inches(0.45))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_GREEN
    arrow.line.color.rgb = ACCENT_GREEN

# ── Slide: AgenticScope ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
            "AgenticScope: Unified State Management", font_size=36, color=TITLE_COLOR,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.45),
            "All 11 patterns coordinate through one shared AgenticScope",
            font_size=19, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.9), Inches(1.95), Inches(11.5), Inches(1.25),
            "•  Shared state via readState() / writeState()     •  Output mapping via @Agent(outputKey = ...)\n"
            "•  Invocation tracking for debugging              •  AgentListener events for real-time visualization",
            font_size=16, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)

add_panel(slide, Inches(0.75), Inches(3.2), Inches(5.8), Inches(1.65))
add_textbox(slide, Inches(1.05), Inches(3.48), Inches(5.2), Inches(0.4),
            "PROGRAMMATIC", font_size=18, color=TITLE_COLOR, bold=True)
add_textbox(slide, Inches(1.05), Inches(3.95), Inches(5.15), Inches(0.65),
            "Builder APIs such as sequenceBuilder() and loopBuilder()\nBest for dynamic workflows and runtime configuration",
            font_size=14, color=TEXT_COLOR)

add_panel(slide, Inches(6.78), Inches(3.2), Inches(5.8), Inches(1.65), ACCENT_GREEN)
add_textbox(slide, Inches(7.08), Inches(3.48), Inches(5.2), Inches(0.4),
            "DECLARATIVE", font_size=18, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(7.08), Inches(3.95), Inches(5.15), Inches(0.65),
            "Annotations plus createAgenticSystem()\nBest for concise definitions with compile-time validation",
            font_size=14, color=TEXT_COLOR)

flow_items = [
    (Inches(1.05), "Agent A\nwrites output"),
    (Inches(5.1), "AgenticScope\nshared state"),
    (Inches(9.15), "Agent B\nreads input"),
]
for left, label in flow_items:
    flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, Inches(5.6), Inches(3.1), Inches(1.05))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = TABLE_HEADER_BG
    flow_box.line.color.rgb = ACCENT_GREEN if "AgenticScope" in label else TITLE_COLOR
    add_textbox(slide, left, Inches(5.78), Inches(3.1), Inches(0.7), label,
                font_size=15, color=TEXT_COLOR, bold="AgenticScope" in label,
                alignment=PP_ALIGN.CENTER)

for left in (Inches(4.35), Inches(8.4)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                   left, Inches(5.88), Inches(0.55), Inches(0.48))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_GREEN
    arrow.line.color.rgb = ACCENT_GREEN

# ── Save ─────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "MatrixAgents-Showcase.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
